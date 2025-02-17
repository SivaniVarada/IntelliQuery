# ===========================================
#  IntelliQuery
#  A Streamlit application integrating various document retrieval and Q&A functionalities.
# ===========================================

# ===========================================
# 1) Imports and Dependencies
# ===========================================
import os
import re
import shutil
import time
import base64
import tempfile
import ffmpeg
import cohere
import cv2
import google.generativeai as genai
import pandas as pd
import requests
import speech_recognition as sr
import streamlit as st
import torch
import torch._classes
import whisper
from bs4 import BeautifulSoup
from datetime import datetime
from dotenv import load_dotenv
from fpdf import FPDF
from googletrans import Translator, LANGUAGES
from langchain.chains import RetrievalQA
from langchain.chains.question_answering import load_qa_chain
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.prompts import PromptTemplate
from langchain_experimental.agents import create_csv_agent
from langchain_core.documents import Document
from langchain_community.vectorstores import FAISS
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PyPDF2 import PdfReader
from transformers import BertTokenizer, BertForSequenceClassification

# ===========================================
# 2) Global Configurations and Session State
# ===========================================

# Load environment variables from .env file
load_dotenv()

# Initialize Cohere client
cohere_client = cohere.Client(os.getenv("COHERE_API_KEY"))

# Directory to store uploads
UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Initialize Google Generative AI embeddings
embeddings = GoogleGenerativeAIEmbeddings(
    model="models/embedding-001",
    api_key=os.getenv("GOOGLE_API_KEY"),
)

# Initialize or reset session states
if 'conversation_history' not in st.session_state:
    st.session_state.conversation_history = []
if 'content' not in st.session_state:
    st.session_state.content = ""
if 'processing' not in st.session_state:
    st.session_state.processing = False
if 'current_question' not in st.session_state:
    st.session_state.current_question = None
if 'uploaded_image' not in st.session_state:
    st.session_state.uploaded_image = None
if 'uploaded_audio' not in st.session_state:
    st.session_state.uploaded_audio = None
if 'uploaded_video' not in st.session_state:
    st.session_state.uploaded_video = None
if 'vector_store' not in st.session_state:
    st.session_state.vector_store = None
if 'documents_processed' not in st.session_state:
    st.session_state.documents_processed = False
if 'audio_processed' not in st.session_state:
    st.session_state.audio_processed = False
if "conversation_history" not in st.session_state:
    st.session_state.conversation_history = []
if "content" not in st.session_state:
    st.session_state.content = ""

# Page-wide Streamlit config
st.set_page_config(layout="wide", page_title="IntelliQuery")

# ===========================================
# 3) Utility Functions for File Handling
# ===========================================

def clear_old_index():
    """
    Delete any old FAISS vector store index if it exists.
    """
    index_folder = "vector_store_index"
    if os.path.exists(index_folder):
        shutil.rmtree(index_folder)
        print("[INFO] Old FAISS index deleted successfully.")

def handle_file_upload(uploaded_file):
    """
    Saves an uploaded file to the UPLOAD_FOLDER and returns the file path.
    """
    if not uploaded_file:
        return None

    os.makedirs(UPLOAD_FOLDER, exist_ok=True)
    file_path = os.path.join(UPLOAD_FOLDER, uploaded_file.name)

    with open(file_path, "wb") as f:
        f.write(uploaded_file.getbuffer())

    print(f"[INFO] File saved at: {file_path}")
    return file_path

def load_excel_and_convert_to_csv(file):
    """
    Read an Excel file and convert it to a CSV format.
    Returns a string of CSV data.
    """
    try:
        df = pd.read_excel(file, engine="openpyxl")
        print(f"[DEBUG] Loaded Excel with {df.shape[0]} rows and {df.shape[1]} columns.")
        return df.to_csv(index=False)
    except Exception as e:
        print(f"[ERROR] Reading Excel file: {e}")
        return f"Error: {str(e)}"

def get_ppt_content(file):
    """
    Extract text content from each slide of a PPTX file.
    """
    slides_content = []
    prs = Presentation(file)
    for slide in prs.slides:
        slide_text = "".join(
            [shape.text for shape in slide.shapes if hasattr(shape, "text")]
        )
        slides_content.append(slide_text)

    print(f"[DEBUG] Extracted {len(slides_content)} slides from PPT.")
    return "\n".join(slides_content)

def get_pdf_text(file):
    """
    Extract all text from a PDF file.
    """
    text = ""
    pdf_reader = PdfReader(file)
    for page in pdf_reader.pages:
        text += page.extract_text()

    print(f"[DEBUG] Extracted {len(text.split())} words from PDF.")
    return text

# ===========================================
# 4) Audio & Video Processing Functions
# ===========================================

def extract_audio(video_path):
    """
    Extract audio from a video file and save it as a WAV file.
    Returns the path to the WAV file or None on failure.
    """
    if not os.path.exists(video_path):
        print(f"[ERROR] Video file does not exist: {video_path}")
        return None

    audio_output = os.path.splitext(video_path)[0] + ".wav"
    try:
        (
            ffmpeg
            .input(video_path)
            .output(audio_output, format='wav', acodec='pcm_s16le', ar='16000')
            .run(overwrite_output=True)
        )
        if os.path.exists(audio_output):
            print(f"[INFO] Audio extracted and saved to: {audio_output}")
            return audio_output
        else:
            print("[ERROR] Audio extraction failed.")
            return None
    except Exception as e:
        print(f"[ERROR] Extracting audio: {e}")
        return None

# Load a Whisper model for audio transcription
WHISPER_MODEL = whisper.load_model("base")

def transcribe_audio(audio_file_path):
    """
    Transcribe speech from an audio file using Whisper.
    Returns the transcription text.
    """
    if not isinstance(audio_file_path, str) or not os.path.exists(audio_file_path):
        print(f"[ERROR] Invalid file path: {audio_file_path}")
        return None

    try:
        result = WHISPER_MODEL.transcribe(audio_file_path)
        transcription = result["text"]
        print(f"[DEBUG] Transcription preview: {transcription[:100]}...")
        return transcription
    except Exception as e:
        print(f"[ERROR] Whisper Transcription Error: {e}")
        return None

def process_video(uploaded_video):
    """
    Process an uploaded video file:
    1. Upload and save the video locally.
    2. Extract audio from the video.
    3. Transcribe the audio using Whisper.
    4. Update the Vector Store with the transcription.
    5. Clean up temporary files.
    """
    if not uploaded_video:
        return False

    try:
        # Save the video file
        video_path = handle_file_upload(uploaded_video)
        if not video_path:
            return False

        # Extract audio
        audio_path = extract_audio(video_path)
        if not audio_path:
            print("[ERROR] Failed to extract audio from video.")
            return False

        # Transcribe the audio
        transcription = transcribe_audio(audio_path)
        if transcription:
            # Update session content
            if 'content' not in st.session_state:
                st.session_state.content = ""
            st.session_state.content += transcription + "\n"

            # Rebuild the vector store
            clear_old_index()
            text_chunks = get_text_chunks(st.session_state.content)
            st.session_state.vector_store = get_vector_store(text_chunks, "vector_store_index")
            st.session_state.documents_processed = True

            # Clean up temporary files
            try:
                os.remove(video_path)
                os.remove(audio_path)
                print("[INFO] Temporary files cleaned up.")
            except Exception as e:
                print(f"[WARNING] Could not remove temporary files: {e}")

            return True
        return False

    except Exception as e:
        print(f"[ERROR] Processing video: {e}")
        return False

def process_audio(uploaded_audio):
    """
    Process an uploaded audio file:
    1. Upload and save the audio locally.
    2. Transcribe using Whisper.
    3. Create or update the Vector Store with optimized chunks.
    4. Clean up the temporary file.
    """
    if not uploaded_audio:
        return False

    try:
        # Save the audio file
        audio_path = handle_file_upload(uploaded_audio)
        if not audio_path:
            return False

        # Transcribe the audio
        transcription = transcribe_audio(audio_path)
        if transcription:
            # Clear old index
            clear_old_index()

            # Break transcription into smaller chunks for better retrieval
            text_chunks = RecursiveCharacterTextSplitter(
                chunk_size=500,
                chunk_overlap=50,
                separators=["\n\n", "\n", ".", "!", "?", ",", " "]
            ).split_text(transcription)

            # Create the FAISS vector store
            st.session_state.vector_store = FAISS.from_texts(
                texts=text_chunks,
                embedding=embeddings,
                normalize_L2=True
            )

            # Save the updated vector store locally
            st.session_state.vector_store.save_local("vector_store_index")
            st.session_state.documents_processed = True

            # Update session content
            if 'content' not in st.session_state:
                st.session_state.content = ""
            st.session_state.content = transcription  # Replace existing content

            # Clean up temporary audio file
            try:
                os.remove(audio_path)
                print(f"[INFO] Temporary file removed: {audio_path}")
            except Exception as e:
                print(f"[WARNING] Could not remove file: {e}")

            return True
        return False

    except Exception as e:
        print(f"[ERROR] Processing audio: {e}")
        return False

# ===========================================
# 5) Text Processing & Vector Store Functions
# ===========================================

def get_text_chunks(text):
    """
    Split the provided text into chunks for vector storage.
    """
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
    return text_splitter.split_text(text)

def clean_text(text):
    """
    Remove non-UTF-8 characters, emojis, and surrogate Unicode pairs to ensure
    text is clean before embedding.
    """
    if not text:
        return ""

    # Replace invalid UTF-8 chars
    text = text.encode("utf-8", "replace").decode("utf-8")

    # Remove emojis & special symbols
    text = re.sub(r'[\U00010000-\U0010ffff]', '', text)

    return text.strip()

def get_vector_store(text_chunks, vector_store_path):
    """
    Create a FAISS vector store from text chunks after cleaning.
    """
    cleaned_chunks = [clean_text(chunk) for chunk in text_chunks]
    vector_store = FAISS.from_texts(cleaned_chunks, embedding=embeddings)
    vector_store.save_local(vector_store_path)
    print(f"[INFO] FAISS index created with {len(cleaned_chunks)} chunks.")
    return vector_store

def load_vector_store(vector_store_path):
    """
    Load a FAISS vector store from the specified path, or return None if it doesn't exist.
    """
    if not os.path.exists(vector_store_path):
        print("[WARNING] No FAISS index found. Returning an empty vector store.")
        return None
    print(f"[INFO] FAISS index loaded from: {vector_store_path}")
    return FAISS.load_local(vector_store_path, embeddings, allow_dangerous_deserialization=True)

# ===========================================
# 6) Q&A Chain and Retrieval Functions
# ===========================================

def get_conversational_chain():
    """
    Return a Q&A chain configured for a ChatGoogleGenerativeAI model (Gemini 1.5).
    """
    prompt_template = """
    Answer the question as detailed as possible from the provided context. 
    If the answer is not in the context, respond with: "Answer is not available in the context."
    Context:\n {context}\nQuestion:\n{question}\nAnswer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-1.5-flash-latest", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    return load_qa_chain(model, chain_type="stuff", prompt=prompt)

def get_late_chunked_text(retrieved_docs, chunk_size=1000, chunk_overlap=100):
    """
    Dynamically chunk retrieved documents to maintain structure in the final response.
    """
    chunked_docs = []

    for doc in retrieved_docs:
        if isinstance(doc, Document):
            text = clean_text(doc.page_content)
        elif isinstance(doc, dict) and "page_content" in doc:
            text = clean_text(doc["page_content"])
        else:
            text = clean_text(str(doc))

        start = 0
        while start < len(text):
            chunked_docs.append(Document(page_content=text[start : start + chunk_size]))
            start += chunk_size - chunk_overlap

    return chunked_docs

def retrieve_documents(query):
    """
    Search the available vector store using the query. If no store is available, returns an empty list.
    """
    try:
        if st.session_state.vector_store is not None:
            vector_store = st.session_state.vector_store
        else:
            vector_store = load_vector_store("vector_store_index")
            if vector_store:
                st.session_state.vector_store = vector_store

        if not vector_store:
            print("[WARNING] No vector store available for retrieval.")
            return []

        retrieved_docs = vector_store.similarity_search(query)

        # Debug info
        print("\n==============================")
        print(f"[DEBUG] Query: {query}")
        print(f"[DEBUG] Retrieved {len(retrieved_docs)} relevant documents.")
        for i, doc in enumerate(retrieved_docs[:3]):
            print(f"[DEBUG] Document {i+1} preview: {doc.page_content[:200]}...")
        print("==============================\n")

        return retrieved_docs

    except Exception as e:
        print(f"[ERROR] Retrieving documents: {e}")
        return []

def process_question(question, retrieved_docs):
    """
    Process the user's question using retrieved documents and a Q&A chain.
    """
    if not retrieved_docs:
        print(f"[WARNING] No documents found for query: {question}")
        return "No relevant documents found. Please upload documents first or try a different question."

    chain = get_conversational_chain()
    response = chain({"input_documents": retrieved_docs, "question": question}, return_only_outputs=True)

    print(f"[DEBUG] Response generated preview: {response['output_text'][:200]}...")
    return response["output_text"]

# ===========================================
# 7) Query Improvement and Processing
# ===========================================

def fetch_related_terms(query):
    """
    Use Cohere to fetch a list of related terms for short queries, improving retrieval.
    """
    try:
        response = cohere_client.generate(
            model="command",
            prompt=(
                f"Provide a list of related search terms (separated by commas) for improving retrieval. "
                f"Do NOT change the meaning of the query: '{query}'"
            ),
            max_tokens=15
        )

        related_terms = response.generations[0].text.strip()
        # Clean and join the terms
        related_terms = ", ".join([term.strip() for term in related_terms.split(",") if term.strip()])

        print("\n==============================")
        print(f"[DEBUG] Original Query: {query}")
        print(f"[DEBUG] Related Terms: {related_terms}")
        print("==============================\n")

        return related_terms
    except Exception as e:
        print(f"[ERROR] Fetching related terms from Cohere: {e}")
        return ""

def process_input(question):
    """
    Decide how to process the question based on its length.
    - For short queries, fetch related terms to improve retrieval.
    - For longer queries, skip the related terms step.
    """
    if len(question.split()) < 15:
        related_terms = fetch_related_terms(question)
        combined_query = f"{question} {related_terms}" if related_terms else question
        retrieved_docs = retrieve_documents(combined_query)
        return process_question(question, retrieved_docs)
    else:
        retrieved_docs = retrieve_documents(question)
        return process_question(question, retrieved_docs)

def get_gemini_response(question, image):
    """
    Generate a response from the Gemini model for an image-based query.
    """
    model = genai.GenerativeModel('gemini-1.5-flash-latest')
    try:
        response = model.generate_content([question, image])
        return response.text
    except Exception as e:
        return f"Error processing image: {str(e)}"

# ===========================================
# 8) Streamlit UI and Interaction
# ===========================================

def handle_submit():
    """
    Handle the 'Send' button click and process user input.
    """
    if st.session_state.user_input and st.session_state.user_input.strip():
        if not (st.session_state.content or st.session_state.uploaded_image):
            st.error("Please upload the documents before asking questions.")
            return

        st.session_state.current_question = st.session_state.user_input
        st.session_state.processing = True
        st.session_state.user_input = ""

# ===========================================
# 9) Generating Conversation PDF
# ===========================================

class ChatPDF(FPDF):
    """
    Custom PDF class to format and store the conversation history.
    """

    def __init__(self, file_names):
        super().__init__()
        self.file_names = file_names
        self.set_auto_page_break(auto=True, margin=15)
        self.add_page()
        self.set_margins(20, 20, 20)

    def header(self):
        """
        Custom header with title, file references, and generation timestamp.
        """
        self.set_font('helvetica', 'B', 16)
        self.cell(0, 10, 'IntelliQuery Conversation', 0, 1, 'C')
        self.ln(5)

        if self.file_names:
            self.set_font('helvetica', 'I', 10)
            file_names_str = ", ".join(self.file_names)
            self.cell(0, 10, f'Files: {file_names_str}', 0, 1, 'L')

        self.set_font('helvetica', 'I', 10)
        self.cell(0, 10, f'Generated on: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}', 0, 1, 'R')
        self.line(10, self.get_y(), self.w - 10, self.get_y())
        self.ln(10)

    def footer(self):
        """
        Custom footer with hyperlink and page number.
        """
        self.set_y(-15)
        self.set_font('helvetica', 'I', 8)

        # Hyperlink (left)
        self.set_x(10)
        self.set_text_color(0, 102, 204)
        self.cell(0, 10, 'IntelliQuery built by Pranav Vuddagiri', link="https://example.com", align='L')

        # Page number (right)
        self.set_x(-30)
        self.set_text_color(0)
        self.cell(0, 10, f'Page {self.page_no()}', align='R')

    def add_message(self, role, content):
        """
        Add a formatted role (User/Assistant) and content to the PDF.
        """
        initial_y = self.get_y()
        self.line(self.l_margin, initial_y - 2, self.w - self.r_margin, initial_y - 2)
        self.ln(5)

        message_start_y = self.get_y()
        self.set_font('helvetica', 'B', 12)
        role_label_width = 30
        content_x_start = self.l_margin + role_label_width
        content_width = self.w - self.r_margin - content_x_start

        self.set_xy(self.l_margin, message_start_y)
        self.cell(role_label_width, 10, f"{role}:", 0, 0, 'L')

        self.set_font('helvetica', '', 11)
        for line in content.splitlines():
            self.set_xy(content_x_start, message_start_y)
            self.multi_cell(content_width, 10, line, 0, 'J')
            message_start_y = self.get_y()

        self.ln(5)

def create_download_pdf(file_names):
    """
    Generate a PDF of the conversation history and return it as bytes.
    """
    try:
        pdf = ChatPDF(file_names)
        pdf.alias_nb_pages()

        if "conversation_history" in st.session_state and st.session_state.conversation_history:
            for question, answer in st.session_state.conversation_history:
                pdf.add_message("User", question)
                pdf.add_message("Assistant", answer)
        else:
            pdf.add_message("System", "No conversations available.")

        return pdf.output(dest='S').encode('latin1')
    except Exception as e:
        st.error(f"Error generating PDF: {str(e)}")
        return None

# ===========================================
# 10) Streamlit Page Layout
# ===========================================

# Sidebar: Logo, Download Button, and File Uploader
with st.sidebar:
    st.image("logo.svg", width=250)

    # Custom download button styling
    st.markdown("""
        <style>
        .stDownloadButton {
            background-color: transparent !important;
            color: #fff !important;
            border: 1px solid #262730 !important;
            padding: 10px 10px !important;
            border-radius: 4px !important;
            width: 20% !important;
            margin: 5px auto !important;
            display: block !important;
            text-align: center !important;
            font-weight: 500 !important;
        }
        .stDownloadButton:hover {
            background-color: #000 !important;
        }
        </style>
    """, unsafe_allow_html=True)

    # PDF download button
    st.markdown('<div style="text-align: center;">', unsafe_allow_html=True)
    pdf_data = create_download_pdf(file_names=[])
    st.download_button(
        label="📥 Download Conversation",
        data=pdf_data if pdf_data else b"",
        file_name=f"IntelliQuery Conversation {datetime.now()}.pdf",
        mime="application/pdf"
    )
    st.markdown('</div>', unsafe_allow_html=True)

    st.title("Upload Your Documents")
    file_type = st.selectbox("Select file type", ["PDF", "PPT", "Excel", "Image", "Audio", "Video"])
    uploaded_files = None
    uploaded_audio = None
    uploaded_video = None

    # Handle PDF Upload
    if file_type == "PDF":
        uploaded_files = st.file_uploader("Choose PDF files", type=["pdf"], accept_multiple_files=True)

    # Handle Audio Upload
    elif file_type == "Audio":
        uploaded_audio = st.file_uploader("Choose an audio file", type=["mp3", "wav"], key="audio_uploader")
        if uploaded_audio is not None and not st.session_state.audio_processed:
            with st.spinner("Processing audio file..."):
                if process_audio(uploaded_audio):
                    st.session_state.audio_processed = True
                    st.success("✅ Audio processed successfully! You can now ask questions about its content.")
                else:
                    st.error("❌ Error processing the audio file.")

    # Handle Video Upload
    elif file_type == "Video":
        if 'video_processed' not in st.session_state:
            st.session_state.video_processed = False
        uploaded_video = st.file_uploader("Choose a video file", type=["mp4", "avi"], key="video_uploader")
        if uploaded_video is not None and not st.session_state.video_processed:
            with st.spinner("Processing video file..."):
                if process_video(uploaded_video):
                    st.session_state.video_processed = True
                    st.success("✅ Video processed successfully! You can now ask questions about its content.")
                else:
                    st.error("❌ Error processing the video file.")

    # Handle PPT Upload
    elif file_type == "PPT":
        uploaded_files = st.file_uploader("Choose PPT files", type=["pptx"], accept_multiple_files=True)

    # Handle Excel Upload
    elif file_type == "Excel":
        uploaded_files = st.file_uploader("Choose Excel files", type=["xlsx"], accept_multiple_files=True)

    # Handle Image Upload
    elif file_type == "Image":
        uploaded_file = st.file_uploader("Upload an image", type=["jpg", "jpeg", "png"])
        if uploaded_file is not None:
            image = Image.open(uploaded_file)
            st.session_state.uploaded_image = image
            st.image(image, caption="Uploaded Image", use_container_width=True)
            st.success("Image uploaded successfully! You can now ask questions about the image.")

    # Process any uploaded PDF/PPT/Excel files
    if file_type != "Image" and uploaded_files:
        combined_content = ""

        if not st.session_state.documents_processed:
            # Clear existing FAISS index before processing new files
            clear_old_index()

            for file in uploaded_files:
                if file_type == "PDF":
                    combined_content += get_pdf_text(file) + "\n"
                elif file_type == "PPT":
                    combined_content += get_ppt_content(file) + "\n"
                elif file_type == "Excel":
                    combined_content += load_excel_and_convert_to_csv(file) + "\n"

            # Update session content
            st.session_state.content = combined_content

            # Create FAISS index from the combined content
            text_chunks = get_text_chunks(combined_content)
            st.session_state.vector_store = get_vector_store(text_chunks, "vector_store_index")
            st.session_state.documents_processed = True

            st.success(f"✅ {len(uploaded_files)} files processed successfully!")

# ===========================================
# 11) Main Page Layout
# ===========================================
header_container = st.container()
with header_container:
    st.markdown("<h1>IntelliQuery: Empowering Precision with RAG</h1>", unsafe_allow_html=True)

# Collect file names for PDF creation
if "uploaded_files" in st.session_state:
    file_names = [file.name for file in st.session_state.uploaded_files]
else:
    file_names = []

# Additional CSS for the download button
st.markdown("""
<style>
    .stDownloadButton {
        position: fixed;
        top: 20px;
        right: 20px;
        z-index: 999;
        background-color: #000000;
        color: white;
        border-radius: 4px;
        padding: 0.5rem 1rem;
        font-weight: bold;
    }
    .stDownloadButton:hover {
        background-color: #0A6AAE;
        color: white;
    }
</style>
""", unsafe_allow_html=True)

# ===========================================
# 12) Chat Display and Logic
# ===========================================
chat_placeholder = st.container()
with chat_placeholder:
    # Display existing conversation
    for q, a in st.session_state.conversation_history:
        st.markdown(f"""
            <div class='message-container'>
                <div class='user-message-container'>
                    <div class='user-message'>
                        <div class='message-icon'>👤</div>
                        <div class='message-content'>{q}</div>
                    </div>
                </div>
                <div class='bot-message-container'>
                    <div class='bot-message'>
                        <div class='message-icon'>🤖</div>
                        <div class='message-content'>{a}</div>
                    </div>
                </div>
            </div>
        """, unsafe_allow_html=True)

    # Display loading indicator if a question is being processed
    if st.session_state.processing and st.session_state.current_question:
        st.markdown(f"""
            <div class='message-container'>
                <div class='user-message-container'>
                    <div class='user-message'>
                        <div class='message-icon'>👤</div>
                        <div class='message-content'>{st.session_state.current_question}</div>
                    </div>
                </div>
                <div class='bot-message-container'>
                    <div class='bot-message loading'>
                        <div class='dots-container'>
                            <div class='dot'></div>
                            <div class='dot'></div>
                            <div class='dot'></div>
                        </div>
                    </div>
                </div>
            </div>
        """, unsafe_allow_html=True)

        # Determine whether to process an image or text-based question
        if st.session_state.uploaded_image is not None:
            response = get_gemini_response(st.session_state.current_question, st.session_state.uploaded_image)
        else:
            response = process_input(st.session_state.current_question)

        # Update conversation history with the new response
        st.session_state.conversation_history.append((st.session_state.current_question, response))
        st.session_state.processing = False
        st.session_state.current_question = None
        st.rerun()

# ===========================================
# 13) Input Box for Questions
# ===========================================
input_container = st.container()
with input_container:
    col1, col2 = st.columns([6, 1])
    with col1:
        st.text_input(
            label="",
            placeholder="Type your question here...",
            key="user_input",
            on_change=handle_submit,
            label_visibility="collapsed"
        )
    with col2:
        st.button("Send", on_click=handle_submit)

# ===========================================
# 14) Additional CSS
# ===========================================
st.markdown("""
<style>
/* Main container spacing */
.main .block-container {
    padding-bottom: 100px !important;
}

/* Message container styles */
.message-container {
    margin: 1rem 0;
    padding: 0 1rem;
}

/* Loading animation styling for bot messages */
.bot-message.loading {
    background-color: #1A1A1A;
    padding: 1rem;
}

.dots-container {
    display: flex;
    gap: 4px;
}

.dot {
    width: 8px;
    height: 8px;
    background-color: #0E86D4;
    border-radius: 50%;
    animation: bounce 1.4s infinite ease-in-out;
}

.dot:nth-child(1) {
    animation-delay: -0.32s;
}
.dot:nth-child(2) {
    animation-delay: -0.16s;
}

@keyframes bounce {
    0%, 80%, 100% {
        transform: translateY(0);
    }
    40% {
        transform: translateY(-8px);
    }
}

/* Fixed input styling */
.stTextInput input {
    position: fixed !important;
    bottom: 20px !important;
    left: calc(370px + 0.5rem) !important;
    width: calc(100% - (370px + 0.5rem + 8rem)) !important;
    background: #000 !important;
    padding: 15px 20px !important;
    border-radius: 8px !important;
    border: 1px solid #4A4A4A !important;
    color: white !important;
    font-size: 0.9rem !important;
    height: 40px !important;
    z-index: 999 !important;
}

/* 'Send' button styling */
.stButton button {
    position: fixed !important;
    bottom: 20px !important;
    right: 2rem !important;
    background: transparent !important;
    color: #0E86D4 !important;
    width: 80px !important;
    height: 40px !important;
    border-radius: 6px !important;
    cursor: pointer !important;
    z-index: 999 !important;
}

/* Hide default Streamlit elements */
[data-testid="stToolbar"],
[data-testid="stDecoration"],
[data-testid="stStatusWidget"],
footer {
    display: none !important;
}

/* User/bot message styling */
.user-message-container {
    display: flex;
    justify-content: flex-end;
    margin-bottom: 0.5rem;
}
.bot-message-container {
    display: flex;
    justify-content: flex-start;
    margin-bottom: 0.5rem;
}
.user-message, .bot-message {
    padding: 0.8rem;
    border-radius: 15px;
    max-width: 80%;
    display: flex;
    align-items: flex-start;
}
.user-message {
    background-color: #2C3333;
    border-radius: 15px 15px 0 15px;
    margin-left: auto;
}
.bot-message {
    background-color: #1A1A1A;
    border-radius: 15px 15px 15px 0;
}
.message-icon {
    margin-right: 0.5rem;
    font-size: 1.2rem;
}
.message-content {
    color: #FFFFFF;
    line-height: 1.5;
}

/* Mobile responsiveness */
@media (max-width: 768px) {
    .stTextInput input {
        left: 1rem !important;
        width: calc(100% - 7rem) !important;
        bottom: 10px !important;
    }
    .stButton button {
        right: 0.5rem !important;
        bottom: 10px !important;
    }
}
</style>
""", unsafe_allow_html=True)
